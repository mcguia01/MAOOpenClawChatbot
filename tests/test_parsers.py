"""Unit tests for document parsers (docx, pptx, xlsx).

Fixtures are created in-memory using the respective libraries so no
external files are needed on disk.
"""

import io
import tempfile
from pathlib import Path

import docx
import openpyxl
import pytest
from pptx import Presentation
from pptx.util import Inches

from ingestion.parsers.docx_parser import parse as parse_docx
from ingestion.parsers.pptx_parser import parse as parse_pptx
from ingestion.parsers.xlsx_parser import parse as parse_xlsx


# ── Fixtures ──────────────────────────────────────────────────────────────────


@pytest.fixture()
def sample_docx(tmp_path: Path) -> Path:
    """Create a minimal .docx file with paragraphs and a table."""
    doc = docx.Document()
    doc.add_heading("MAO Order Process", level=1)
    doc.add_paragraph("This document describes the MAO order process.")
    doc.add_paragraph("Orders must be confirmed within 24 hours.")

    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Field"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Status"
    table.cell(1, 1).text = "Active"

    file_path = tmp_path / "test.docx"
    doc.save(str(file_path))
    return file_path


@pytest.fixture()
def sample_pptx(tmp_path: Path) -> Path:
    """Create a minimal .pptx file with two slides."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[1]  # Title and Content layout

    slide1 = prs.slides.add_slide(blank_layout)
    slide1.shapes.title.text = "Introduction"
    slide1.placeholders[1].text = "Overview of MAO processes."

    slide2 = prs.slides.add_slide(blank_layout)
    slide2.shapes.title.text = "Order Flow"
    slide2.placeholders[1].text = "Orders flow from MAO to MAWM."

    # Add speaker notes to slide 2
    notes_slide = slide2.notes_slide
    notes_slide.notes_text_frame.text = "Remind audience to check the SOP."

    file_path = tmp_path / "test.pptx"
    prs.save(str(file_path))
    return file_path


@pytest.fixture()
def sample_xlsx(tmp_path: Path) -> Path:
    """Create a minimal .xlsx file with a header row and data rows."""
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Orders"
    ws.append(["OrderNumber", "Status", "Facility"])
    ws.append(["ORD001", "Allocated", "DC01"])
    ws.append(["ORD002", "Shipped", "DC02"])

    file_path = tmp_path / "test.xlsx"
    wb.save(str(file_path))
    return file_path


# ── docx parser tests ──────────────────────────────────────────────────────────


class TestDocxParser:
    def test_returns_list_of_dicts(self, sample_docx: Path) -> None:
        result = parse_docx(sample_docx)
        assert isinstance(result, list)
        assert len(result) > 0

    def test_paragraph_chunk_has_required_keys(self, sample_docx: Path) -> None:
        result = parse_docx(sample_docx)
        para_chunks = [c for c in result if c["section"] != "table"]
        assert len(para_chunks) > 0
        chunk = para_chunks[0]
        assert "text" in chunk
        assert "source" in chunk
        assert "section" in chunk
        assert "page" in chunk

    def test_table_chunks_extracted(self, sample_docx: Path) -> None:
        result = parse_docx(sample_docx)
        table_chunks = [c for c in result if c["section"] == "table"]
        assert len(table_chunks) > 0
        # Table rows should be pipe-delimited
        assert "|" in table_chunks[0]["text"]

    def test_source_is_filename(self, sample_docx: Path) -> None:
        result = parse_docx(sample_docx)
        assert all(c["source"] == "test.docx" for c in result)

    def test_page_is_none(self, sample_docx: Path) -> None:
        result = parse_docx(sample_docx)
        assert all(c["page"] is None for c in result)


# ── pptx parser tests ─────────────────────────────────────────────────────────


class TestPptxParser:
    def test_returns_list_of_dicts(self, sample_pptx: Path) -> None:
        result = parse_pptx(sample_pptx)
        assert isinstance(result, list)
        assert len(result) > 0

    def test_slide_section_label(self, sample_pptx: Path) -> None:
        result = parse_pptx(sample_pptx)
        sections = [c["section"] for c in result]
        assert any("Slide 1" in str(s) for s in sections)
        assert any("Slide 2" in str(s) for s in sections)

    def test_speaker_notes_extracted(self, sample_pptx: Path) -> None:
        result = parse_pptx(sample_pptx)
        notes_chunks = [c for c in result if "Notes" in str(c.get("section", ""))]
        assert len(notes_chunks) > 0
        assert "SOP" in notes_chunks[0]["text"]

    def test_page_is_slide_number(self, sample_pptx: Path) -> None:
        result = parse_pptx(sample_pptx)
        pages = [c["page"] for c in result if "Notes" not in str(c.get("section", ""))]
        assert all(isinstance(p, int) and p >= 1 for p in pages)

    def test_source_is_filename(self, sample_pptx: Path) -> None:
        result = parse_pptx(sample_pptx)
        assert all(c["source"] == "test.pptx" for c in result)


# ── xlsx parser tests ─────────────────────────────────────────────────────────


class TestXlsxParser:
    def test_returns_list_of_dicts(self, sample_xlsx: Path) -> None:
        result = parse_xlsx(sample_xlsx)
        assert isinstance(result, list)
        assert len(result) == 2  # 2 data rows

    def test_row_text_format(self, sample_xlsx: Path) -> None:
        result = parse_xlsx(sample_xlsx)
        first = result[0]["text"]
        assert "OrderNumber: ORD001" in first
        assert "Status: Allocated" in first
        assert "|" in first

    def test_section_is_sheet_name(self, sample_xlsx: Path) -> None:
        result = parse_xlsx(sample_xlsx)
        assert all(c["section"] == "Orders" for c in result)

    def test_page_is_row_number(self, sample_xlsx: Path) -> None:
        result = parse_xlsx(sample_xlsx)
        assert result[0]["page"] == 1
        assert result[1]["page"] == 2

    def test_source_is_filename(self, sample_xlsx: Path) -> None:
        result = parse_xlsx(sample_xlsx)
        assert all(c["source"] == "test.xlsx" for c in result)
