"""Parser for .pptx files using python-pptx.

Extracts:
  - Text from all shapes on each slide
  - Speaker notes (if present)

Each returned dict has the shape:
  {
      "text": str,
      "source": str,        # filename
      "section": str,       # "Slide {n}: {slide_title}" or "Slide {n}: Notes"
      "page": int,          # slide number (1-based)
  }
"""

import logging
from pathlib import Path
from typing import Any

from pptx import Presentation

logger = logging.getLogger(__name__)


def _get_slide_title(slide: Any) -> str:
    """Return the title text of a slide, or an empty string if none."""
    if slide.shapes.title and slide.shapes.title.text:
        return slide.shapes.title.text.strip()
    return ""


def parse(file_path: Path) -> list[dict[str, str | int | None]]:
    """Extract text chunks from a .pptx file.

    Args:
        file_path: Absolute path to the .pptx file.

    Returns:
        List of chunk dicts ready for the chunker.
    """
    prs = Presentation(str(file_path))
    source = file_path.name
    chunks: list[dict[str, str | int | None]] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        title = _get_slide_title(slide)
        section_label = f"Slide {slide_num}: {title}" if title else f"Slide {slide_num}"

        # Extract text from all shapes
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text and text != title:
                    slide_texts.append(text)

        if slide_texts:
            chunks.append(
                {
                    "text": "\n".join(slide_texts),
                    "source": source,
                    "section": section_label,
                    "page": slide_num,
                }
            )

        # Extract speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                chunks.append(
                    {
                        "text": notes_text,
                        "source": source,
                        "section": f"Slide {slide_num}: Notes",
                        "page": slide_num,
                    }
                )

    logger.info("pptx_parser: extracted %d chunks from '%s'", len(chunks), source)
    return chunks
